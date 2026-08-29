"""The three downloadables: the workbook, the chart pack, the slide deck.

Everything is written from a finished run, so a file can never disagree with
what the dashboard showed. The chart palette is the interface's own, so a chart
pasted into a report still looks like the app it came from.

Charts follow the rules that matter for a chart someone will act on: one axis
per chart (never two scales), a category's colour follows the category and not
its rank, values labelled directly rather than left to a gridline, and no
figure invented that the run did not produce.
"""
import io
import os
import zipfile

# AAU's palette, the same values the interface uses.
GREEN = "#0A7A3A"
BRIGHT = "#0FA64F"
DEEP = "#14563A"
RED = "#E0303F"
META = "#63736A"
PAGE = "#F4F6F5"
HAIR = "#E4EAE6"
INK = "#1A1A1A"

# Colour follows the college, never its position in the sorted list -- a
# filter that changes the order must not repaint the bars.
COLLEGE_COLOR = {
    "College of Education, Humanities and Social Sciences": BRIGHT,
    "College of Engineering": GREEN,
    "College of Pharmacy": RED,
    "College of Business": DEEP,
    "College of Law": "#6B8CAE",
    "College of Communication and Media": "#1F8A57",
    "College of Dentistry": "#C98B5E",
    "College of Nursing": "#8FB89E",
}


def _mpl():
    import matplotlib
    matplotlib.use("Agg")
    import matplotlib.pyplot as plt
    plt.rcParams.update({
        "font.family": ["Archivo", "Helvetica Neue", "Helvetica", "DejaVu Sans"],
        "axes.edgecolor": HAIR,
        "axes.labelcolor": META,
        "text.color": INK,
        "xtick.color": META,
        "ytick.color": META,
        "axes.titlesize": 13,
        "axes.titleweight": "bold",
        "figure.dpi": 160,
        "savefig.bbox": "tight",
        "savefig.facecolor": "white",
    })
    return plt


def _frame(ax, title, note=""):
    ax.set_title(title, loc="left", pad=14, color=INK)
    for side in ("top", "right"):
        ax.spines[side].set_visible(False)
    ax.spines["left"].set_color(HAIR)
    ax.spines["bottom"].set_color(HAIR)
    ax.tick_params(length=0)
    if note:
        ax.annotate(note, (0, -0.16), xycoords="axes fraction",
                    fontsize=8.5, color=META, va="top")


def _short(name):
    return (name.replace("College of ", "")
                .replace("Education, Humanities and Social Sciences",
                         "Education & Humanities"))


# The eight colleges in the order the rest of the app uses, so a stacked bar
# stacks the same way twice running.
COLLEGE_ORDER_LOCAL = [
    ("College of Education, Humanities and Social Sciences", BRIGHT),
    ("College of Engineering", GREEN),
    ("College of Pharmacy", RED),
    ("College of Business", DEEP),
    ("College of Law", "#6B8CAE"),
    ("College of Communication and Media", "#1F8A57"),
    ("College of Dentistry", "#C98B5E"),
    ("College of Nursing", "#8FB89E"),
]


def _prog_label(name, room=40):
    """A programme name short enough for an axis, with its degree intact.

    Dropping the degree collapses "Bachelor of Arts in Applied Sociology" and
    "Master of Arts in Applied Sociology" into one label for two programmes
    with different numbers.
    """
    import re
    deg = ""
    for pat, tag in (("Bachelor of Science", "BSc"), ("Master of Science", "MSc"),
                     ("Bachelor of Arts", "BA"), ("Master of Arts", "MA"),
                     ("Doctor of Philosophy", "PhD"),
                     ("Bachelor of Education", "BEd"),
                     ("Master of Education", "MEd"),
                     ("Postgraduate Professional Diploma", "Dip"),
                     ("BBA", "BBA")):
        if name.startswith(pat):
            deg = tag
            break
    t = re.sub(r"^(Bachelor|Master|Doctor) of (Science|Arts|Philosophy|Education)"
               r"( in | - )", "", name)
    t = re.sub(r"^(Bachelor|Master|Doctor) of ", "", t)
    t = re.sub(r"^Postgraduate Professional Diploma in ", "", t)
    t = re.sub(r"^BBA in ", "", t)
    if len(t) > room:
        t = t[:room].rsplit(" ", 1)[0] + "\u2026"
    return ("%s \u00b7 %s" % (deg, t)) if deg else t


def _save(plt, outdir, name):
    p = os.path.join(outdir, name)
    plt.savefig(p)
    plt.close()
    return p


def charts(vm, outdir):
    """Write the chart pack as PNGs. -> [paths]"""
    plt = _mpl()
    os.makedirs(outdir, exist_ok=True)
    made = []

    cols = [c for c in vm["colleges"] if c.get("papers")]
    cols.sort(key=lambda c: c["papers"])

    # 1 -- papers by college. Horizontal because the labels are long; values
    # labelled at the bar end so no one has to read a gridline.
    if cols:
        fig, ax = plt.subplots(figsize=(8.4, 4.6))
        names = [_short(c["name"]) for c in cols]
        vals = [c["papers"] for c in cols]
        colors = [COLLEGE_COLOR.get(c["name"], GREEN) for c in cols]
        bars = ax.barh(names, vals, color=colors, height=.62)
        for b, v in zip(bars, vals):
            ax.annotate("%s" % f"{v:,}", (b.get_width(), b.get_y() + b.get_height() / 2),
                        xytext=(6, 0), textcoords="offset points",
                        va="center", fontsize=10, fontweight="bold", color=INK)
        ax.set_xlim(0, max(vals) * 1.14)
        ax.get_xaxis().set_visible(False)
        _frame(ax, "Papers by college",
               "A paper written across two colleges is credited to both, so "
               "these add up to more than the paper count.")
        p = os.path.join(outdir, "01_papers_by_college.png")
        fig.savefig(p); plt.close(fig); made.append(p)

    # 2 -- people by college, faculty vs everyone else. Two series, stacked,
    # with a 2px white gap so the segments read as separate.
    if cols:
        fig, ax = plt.subplots(figsize=(8.4, 4.6))
        cc = sorted(vm["colleges"], key=lambda c: c.get("people") or 0)
        names = [_short(c["name"]) for c in cc]
        fac = [c.get("people") or 0 for c in cc]
        ax.barh(names, fac, color=GREEN, height=.62, label="On the roster")
        for y, v in enumerate(fac):
            if v:
                ax.annotate(str(v), (v, y), xytext=(6, 0),
                            textcoords="offset points", va="center",
                            fontsize=10, fontweight="bold", color=INK)
        ax.set_xlim(0, max(fac + [1]) * 1.14)
        ax.get_xaxis().set_visible(False)
        _frame(ax, "Faculty on the roster, by college")
        p = os.path.join(outdir, "02_faculty_by_college.png")
        fig.savefig(p); plt.close(fig); made.append(p)

    # 3 -- who the authors are. One bar, split, because the point is the
    # proportion and not the two numbers side by side.
    s = vm["stats"]
    fac, other = s.get("faculty", 0), s.get("authors", 0) - s.get("faculty", 0)
    if fac or other:
        fig, ax = plt.subplots(figsize=(8.4, 1.9))
        ax.barh([""], [fac], color=GREEN, height=.5)
        ax.barh([""], [other], left=[fac + max(1, (fac + other) // 300)],
                color="#C3D6CA", height=.5)
        ax.annotate("%s on the roster" % f"{fac:,}", (fac / 2, 0),
                    ha="center", va="center", color="white",
                    fontsize=11, fontweight="bold")
        ax.annotate("%s students or outside faculty" % f"{other:,}",
                    (fac + other / 2, 0), ha="center", va="center",
                    color=INK, fontsize=11, fontweight="bold")
        ax.set_xlim(0, fac + other)
        ax.axis("off")
        _frame(ax, "Every author on an AAU paper")
        p = os.path.join(outdir, "03_authors_split.png")
        fig.savefig(p); plt.close(fig); made.append(p)

    # 4 -- the most published people on the roster, by AAU papers in the window
    top = [a for a in vm["authors"] if a.get("tag") == "Faculty"]
    top.sort(key=lambda a: -(a.get("papers") or 0))
    top = [a for a in top[:15] if a.get("papers")][::-1]
    if top:
        fig, ax = plt.subplots(figsize=(8.4, 5.4))
        names = [a["name"] for a in top]
        vals = [a["papers"] for a in top]
        bars = ax.barh(names, vals, color=GREEN, height=.62)
        for b, v in zip(bars, vals):
            ax.annotate(str(v), (b.get_width(), b.get_y() + b.get_height() / 2),
                        xytext=(6, 0), textcoords="offset points", va="center",
                        fontsize=10, fontweight="bold", color=INK)
        ax.set_xlim(0, max(vals) * 1.12)
        ax.get_xaxis().set_visible(False)
        ax.tick_params(axis="y", labelsize=9.5)
        _frame(ax, "Most published on the roster",
               "AAU papers in the current window.")
        p = os.path.join(outdir, "04_top_authors.png")
        fig.savefig(p); plt.close(fig); made.append(p)

    # 5 -- h-index against AAU output. Career standing versus what landed in
    # this window; they are different questions and the scatter shows it.
    pts = [(a.get("papers") or 0, a.get("h") or 0, a["name"])
           for a in vm["authors"] if a.get("tag") == "Faculty"
           and (a.get("h") or 0) > 0]
    if len(pts) > 4:
        fig, ax = plt.subplots(figsize=(8.4, 5.0))
        ax.scatter([p[0] for p in pts], [p[1] for p in pts],
                   s=42, color=GREEN, alpha=.72, edgecolor="white", linewidth=.8)
        for x, y, n in sorted(pts, key=lambda t: -(t[0] + t[1]))[:6]:
            ax.annotate(n, (x, y), xytext=(7, 4), textcoords="offset points",
                        fontsize=8.5, color=META)
        ax.set_xlabel("AAU papers in the window", fontsize=9.5)
        ax.set_ylabel("h-index (whole career)", fontsize=9.5)
        ax.grid(axis="y", color=HAIR, linewidth=.8)
        ax.set_axisbelow(True)
        _frame(ax, "Career standing against output in this window")
        p = os.path.join(outdir, "05_h_index_vs_output.png")
        fig.savefig(p); plt.close(fig); made.append(p)

    # ---- 06 · papers per year, stacked by college --------------------------
    # The only temporal view in the pack. The run carries a year on every
    # paper; the view-model's own paper map is capped at fifty rows per author
    # and would undercount every total, so this reads paper_rows.
    rows = vm.get("paper_rows") or []
    if rows:
        years = sorted({r["year"] for r in rows if r.get("year")})
        if len(years) > 1:
            names = [n for n, _ in COLLEGE_ORDER_LOCAL]
            per = {n: [0] * len(years) for n in names}
            per["Not stated"] = [0] * len(years)
            yi = {y: i for i, y in enumerate(years)}
            for r in rows:
                if not r.get("year"):
                    continue
                cs = r.get("college") or []
                for c in (cs or ["Not stated"]):
                    if c in per:
                        per[c][yi[r["year"]]] += 1
            fig, ax = plt.subplots(figsize=(8.4, 4.8))
            bottom = [0] * len(years)
            for n in names + ["Not stated"]:
                v = per[n]
                if not sum(v):
                    continue
                ax.bar([str(y) for y in years], v, bottom=bottom, width=.66,
                       label=_short(n),
                       color=COLLEGE_COLOR.get(n, "#C3D6CA"))
                bottom = [a + b for a, b in zip(bottom, v)]
            _frame(ax, "Papers per year",
                   "A paper with authors from two colleges is counted in both, "
                   "so the bars total more than the papers in the window.")
            ax.legend(frameon=False, fontsize=8, ncol=3, loc="upper left")
            ax.spines["left"].set_visible(False)
            ax.tick_params(axis="y", labelsize=8)
            ax.grid(axis="y", color=HAIR, linewidth=.8)
            ax.set_axisbelow(True)
            made.append(_save(plt, outdir, "06_papers_per_year.png"))

    # ---- 07 · programmes ranked by papers per member of staff --------------
    # The dashboard's whole argument in one chart: size-corrected output, so a
    # three-person programme is comparable with a twenty-one-person one.
    progs = [p for p in (vm.get("programs") or []) if (p.get("tagged") or 0)]
    if len(progs) > 4:
        rank = sorted(progs, key=lambda p: (p.get("papers") or 0) / p["tagged"])
        rank = rank[-18:]
        lab = [_prog_label(p["name"]) for p in rank]
        val = [(p.get("papers") or 0) / p["tagged"] for p in rank]
        col = [COLLEGE_COLOR.get(p["college"], GREEN) for p in rank]
        fig, ax = plt.subplots(figsize=(8.4, 6.2))
        ax.barh(lab, val, height=.62, color=col)
        for i, v in enumerate(val):
            ax.text(v + max(val) * .012, i, "%.1f" % v, va="center",
                    fontsize=8.5, color=META)
        _frame(ax, "Programmes by papers per member of staff",
               "Divided by everyone AAU lists on the programme, whether or not "
               "they have a Scopus record. Colour is the college.")
        ax.set_xlim(0, max(val) * 1.12)
        ax.get_xaxis().set_visible(False)
        ax.tick_params(axis="y", labelsize=8.5)
        made.append(_save(plt, outdir, "07_programmes_per_staff.png"))

    # ---- 08 · impact against volume ---------------------------------------
    # The quadrant a Provost reads first. Medians rather than means, because
    # one 4,000-citation paper would drag a mean across the chart.
    pts = [p for p in (vm.get("programs") or [])
           if (p.get("papers") or 0) > 0 and p.get("citations") is not None]
    if len(pts) > 4:
        x = [p["papers"] for p in pts]
        y = [(p.get("citations") or 0) / p["papers"] for p in pts]
        sz = [28 + 5 * (p.get("tagged") or 1) for p in pts]
        col = [COLLEGE_COLOR.get(p["college"], GREEN) for p in pts]
        fig, ax = plt.subplots(figsize=(8.4, 5.6))
        ax.scatter(x, y, s=sz, c=col, alpha=.75, edgecolor="white", linewidth=.9)
        mx = sorted(x)[len(x) // 2]
        my = sorted(y)[len(y) // 2]
        ax.axvline(mx, color="#B9C4BD", linewidth=1, linestyle=(0, (4, 3)))
        ax.axhline(my, color="#B9C4BD", linewidth=1, linestyle=(0, (4, 3)))
        ax.annotate("median %d papers" % mx, (mx, max(y)),
                    textcoords="offset points", xytext=(5, -2),
                    fontsize=8, color="#8C9A92")
        ax.annotate("median %.1f per paper" % my, (max(x), my),
                    textcoords="offset points", xytext=(-4, 5), ha="right",
                    fontsize=8, color="#8C9A92")
        # Programmes that share a staff list share a coordinate exactly -- five
        # of Communication's do -- so labelling each would stack them into
        # unreadable mush. Label the first at each spot and say how many others
        # sit under it.
        placed = []
        span_x = (max(x) - min(x)) or 1
        span_y = (max(y) - min(y)) or 1
        for pr, xx, yy in sorted(zip(pts, x, y), key=lambda t: -(t[1] * t[2])):
            if len(placed) >= 6:
                break
            near = sum(1 for px, py in placed
                       if abs(px - xx) / span_x < .06
                       and abs(py - yy) / span_y < .06)
            if near:
                continue
            same = sum(1 for qx, qy in zip(x, y)
                       if abs(qx - xx) < 1e-9 and abs(qy - yy) < 1e-9) - 1
            lab_ = _prog_label(pr["name"], 26)
            if same:
                lab_ += "  (+%d alike)" % same
            ax.annotate(lab_, (xx, yy), textcoords="offset points",
                        xytext=(8, 5), fontsize=8, color=META)
            placed.append((xx, yy))
        _frame(ax, "Impact against volume",
               "Each bubble is a programme, sized by the staff AAU lists on it. "
               "Lines are the medians: above and right is more work AND more "
               "cited.")
        ax.set_xlabel("papers in the window", fontsize=9)
        ax.set_ylabel("citations per paper", fontsize=9)
        ax.grid(color=HAIR, linewidth=.7)
        ax.set_axisbelow(True)
        made.append(_save(plt, outdir, "08_impact_against_volume.png"))

    # ---- 09 · the institutions AAU publishes with --------------------------
    net = vm.get("network") or {}
    top = (net.get("top") or [])[:20]
    if top:
        top = list(reversed(top))
        lab = [t["name"][:44] for t in top]
        raw = [t["papers"] for t in top]
        cred = [t.get("credit") or 0 for t in top]
        fig, ax = plt.subplots(figsize=(8.4, 6.4))
        ax.barh(lab, raw, height=.66, color="#CDE6D8")
        ax.barh(lab, cred, height=.66, color=GREEN)
        for i, (r, c) in enumerate(zip(raw, cred)):
            ax.text(r + max(raw) * .012, i, str(r), va="center",
                    fontsize=8.5, color=META)
        _frame(ax, "Who Al Ain University publishes with",
               "Pale is joint papers; solid is the same work shared out among "
               "every institution on each paper. The gap is the consortium.")
        ax.set_xlim(0, max(raw) * 1.1)
        ax.get_xaxis().set_visible(False)
        ax.tick_params(axis="y", labelsize=8.5)
        made.append(_save(plt, outdir, "09_partner_institutions.png"))

    # ---- 10 · where the partners are ---------------------------------------
    ctry = (net.get("countries") or [])[:15]
    if ctry:
        ctry = list(reversed(ctry))
        fig, ax = plt.subplots(figsize=(8.4, 5.2))
        ax.barh([c["name"][:32] for c in ctry], [c["papers"] for c in ctry],
                height=.62, color="#4E9E74")
        mx2 = max(c["papers"] for c in ctry)
        for i, c in enumerate(ctry):
            ax.text(c["papers"] + mx2 * .012, i, str(c["papers"]), va="center",
                    fontsize=8.5, color=META)
        _frame(ax, "Where the partners are",
               "Countries as the addresses print them, counted once per paper.")
        ax.set_xlim(0, mx2 * 1.1)
        ax.get_xaxis().set_visible(False)
        ax.tick_params(axis="y", labelsize=8.5)
        made.append(_save(plt, outdir, "10_partner_countries.png"))

    # ---- 11 · how concentrated the citations are ---------------------------
    # A mean hides this completely. The curve states it.
    cits = sorted((r.get("cited_by") or 0) for r in rows)
    if len(cits) > 20 and sum(cits):
        cits = cits[::-1]
        tot = float(sum(cits))
        cum, run = [], 0.0
        for c in cits:
            run += c
            cum.append(100 * run / tot)
        share = [100 * (i + 1) / len(cits) for i in range(len(cits))]
        fig, ax = plt.subplots(figsize=(8.4, 5.0))
        ax.plot(share, cum, color=GREEN, linewidth=2.2)
        ax.plot([0, 100], [0, 100], color=HAIR, linewidth=1.2, linestyle="--")
        ax.fill_between(share, cum, color=GREEN, alpha=.08)
        half = next((s for s, c in zip(share, cum) if c >= 50), None)
        if half:
            ax.axhline(50, color=HAIR, linewidth=.9)
            ax.annotate("half the citations come from the top %.0f%% of papers"
                        % half, (half, 50), textcoords="offset points",
                        xytext=(8, -16), fontsize=9, color=INK)
        _frame(ax, "How concentrated the citations are",
               "The dashed line is what an even spread would look like. The "
               "further the curve sits above it, the more the citations rest "
               "on a few papers.")
        ax.set_xlabel("share of papers, most cited first (%)", fontsize=9)
        ax.set_ylabel("share of citations (%)", fontsize=9)
        ax.set_xlim(0, 100)
        ax.set_ylim(0, 100)
        ax.grid(color=HAIR, linewidth=.7)
        ax.set_axisbelow(True)
        made.append(_save(plt, outdir, "11_citation_concentration.png"))

    return made


def chart_zip(vm, path):
    """The chart pack as one .zip the browser can download."""
    import tempfile
    tmp = tempfile.mkdtemp(prefix="aau-charts-")
    made = charts(vm, tmp)
    os.makedirs(os.path.dirname(path), exist_ok=True)
    with zipfile.ZipFile(path, "w", zipfile.ZIP_DEFLATED) as z:
        for p in made:
            z.write(p, os.path.basename(p))
    return path, made


def deck(vm, path, generated=""):
    """The slide deck: one title slide, then a slide per chart."""
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    import tempfile

    tmp = tempfile.mkdtemp(prefix="aau-deck-")
    made = charts(vm, tmp)
    s = vm["stats"]

    prs = Presentation()
    prs.slide_width, prs.slide_height = Inches(13.333), Inches(7.5)

    def rgb(h):
        return RGBColor(int(h[1:3], 16), int(h[3:5], 16), int(h[5:7], 16))

    # title
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    bg = sl.shapes.add_shape(1, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid(); bg.fill.fore_color.rgb = rgb(GREEN); bg.line.fill.background()
    tb = sl.shapes.add_textbox(Inches(.9), Inches(2.4), Inches(11.5), Inches(2.6))
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = "Al Ain University research output"
    p.runs[0].font.size = Pt(44); p.runs[0].font.bold = True
    p.runs[0].font.color.rgb = rgb("#FFFFFF")
    p2 = tf.add_paragraph()
    p2.text = ("%s papers · %s authors · %s on the faculty roster"
               % (f"{s.get('papers', 0):,}", f"{s.get('authors', 0):,}",
                  f"{s.get('faculty', 0):,}"))
    p2.runs[0].font.size = Pt(20); p2.runs[0].font.color.rgb = rgb("#D5E8DC")
    p3 = tf.add_paragraph()
    p3.text = ("Nothing is counted unless the paper itself prints an AAU "
               "address. Generated %s." % (generated or "from the latest run"))
    p3.runs[0].font.size = Pt(13); p3.runs[0].font.color.rgb = rgb("#B9D9C6")

    for img in made:
        sl = prs.slides.add_slide(prs.slide_layouts[6])
        sl.shapes.add_picture(img, Inches(.7), Inches(.7), width=Inches(11.9))

    os.makedirs(os.path.dirname(path), exist_ok=True)
    prs.save(path)
    return path
